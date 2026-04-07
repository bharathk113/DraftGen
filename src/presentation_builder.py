from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.util import Inches, Pt
from typing import Optional, List
from document_loader import ExtractedImage

class PresentationBuilder:
    def __init__(self, template_path: Optional[str] = None):
        self.prs = Presentation(template_path) if template_path else Presentation()
        self.default_layout = self._find_layout()

    def _parse_image_id(self, img_id):
        """Parse image ID, handling both int and 'Image X' formats."""
        if isinstance(img_id, int):
            return img_id
        if isinstance(img_id, str):
            if img_id.startswith("Image "):
                try:
                    return int(img_id.split(" ")[1])
                except (ValueError, IndexError):
                    return None
            try:
                return int(img_id)
            except ValueError:
                return None
        return None

    def _find_layout(self):
        if len(self.prs.slide_layouts) > 1:
            return self.prs.slide_layouts[1]  # Usually title and content
        return self.prs.slide_layouts[0]

    def build_from_outline(self, slides, extracted_images: List[ExtractedImage] = None):
        extracted_images = extracted_images or []
        for slide_data in slides:
            title = slide_data.get("title", "Untitled Slide")
            bullets = slide_data.get("bullets", [])
            if isinstance(bullets, str):
                bullets = [b.strip() for b in bullets.split('\n') if b.strip()]
            notes = slide_data.get("notes", "")
            image_ids = slide_data.get("images", [])
            
            slide = self.prs.slides.add_slide(self.default_layout)

            # 1. Set Title Safely
            title_shape = slide.shapes.title
            if title_shape:
                title_shape.text = title

            # 2. Find Content/Body Placeholder
            body = None
            for placeholder in slide.placeholders:
                if placeholder.placeholder_format.type in {PP_PLACEHOLDER.BODY, PP_PLACEHOLDER.OBJECT}:
                    body = placeholder
                    break

            if body is None:
                for shape in slide.shapes:
                    is_title = title_shape and shape.shape_id == title_shape.shape_id
                    if getattr(shape, "has_text_frame", False) and not is_title:
                        body = shape
                        break

            # 3. Add Bullets and Adjust Sizes
            if body and bullets:
                tf = body.text_frame
                
                # -- ENABLE NATIVE POWERPOINT AUTOFIT --
                tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
                tf.word_wrap = True
                
                # -- DYNAMIC FONT SIZING CALCULATION --
                # Count total characters across all bullets in this slide
                total_chars = sum(len(str(b)) for b in bullets)
                
                # Thresholds: Adjust these Pt values if you want it larger or smaller
                if total_chars < 150:
                    calculated_size = Pt(28)
                elif total_chars < 250:
                    calculated_size = Pt(24)
                elif total_chars < 350:
                    calculated_size = Pt(20)
                else:
                    calculated_size = Pt(16)

                # Write first bullet
                p = tf.paragraphs[0]
                p.text = str(bullets[0])
                if p.runs:
                    p.runs[0].font.size = calculated_size
                
                # Write subsequent bullets
                for bullet in bullets[1:]:
                    p = tf.add_paragraph()
                    p.text = str(bullet)
                    p.level = 0
                    if p.runs:
                        p.runs[0].font.size = calculated_size

            # 4. Add Images
            if image_ids and extracted_images:
                # Position images on the right side or below content
                left = Inches(6)  # Right side
                top = Inches(2)
                max_width = Inches(3)
                max_height = Inches(2)
                
                for img_id_raw in image_ids:
                    img_id = self._parse_image_id(img_id_raw)
                    if img_id is not None and 0 <= img_id < len(extracted_images):
                        img = extracted_images[img_id]
                        # Calculate dimensions to preserve aspect ratio
                        aspect_ratio = img.width / img.height
                        if aspect_ratio > max_width / max_height:
                            # Image is wider relative to height
                            width = max_width
                            height = width / aspect_ratio
                        else:
                            # Image is taller relative to width
                            height = max_height
                            width = height * aspect_ratio
                        
                        # Save image temporarily or use in-memory
                        import io
                        img_buffer = io.BytesIO()
                        img.image.save(img_buffer, format='PNG')
                        img_buffer.seek(0)
                        
                        slide.shapes.add_picture(img_buffer, left, top, width, height)
                        top += height + Inches(0.2)  # Stack images vertically

            # 5. Add Notes
            if notes:
                slide.notes_slide.notes_text_frame.text = notes

    def save(self, output_path: str):
        self.prs.save(output_path)