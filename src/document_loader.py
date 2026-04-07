import io
import logging
import os
from pathlib import Path
from typing import Union, Optional, List

from docx import Document as DocxDocument
from PIL import Image
from pptx import Presentation as PptxPresentation
from PyPDF2 import PdfReader


class ExtractedImage:
    """Container for extracted image data with metadata."""
    def __init__(self, image: Image.Image, source: str, page_or_shape: Optional[Union[int, str]] = None):
        self.image = image
        self.source = source  # e.g., "pdf", "docx", "pptx"
        self.page_or_shape = page_or_shape  # page number or shape index
        self.width = image.width
        self.height = image.height
    
    def get_description(self) -> str:
        """Get a brief description of the image."""
        desc = f"Image from {self.source}: {self.width}x{self.height}px"
        if self.page_or_shape:
            desc += f" (page {self.page_or_shape})"
        
        # Try to get OCR text
        try:
            from pytesseract import image_to_string
            text = image_to_string(self.image).strip()
            if text:
                desc += f" - OCR: {text[:100]}{'...' if len(text) > 100 else ''}"
        except ImportError:
            pass
        except Exception:
            pass
        
        return desc


def describe_image(image: Image.Image) -> str:
    """Extract OCR text from image (fallback for text-only mode)."""
    try:
        from pytesseract import image_to_string
        text = image_to_string(image).strip()
        return f"OCR text: {text}" if text else "Image with no readable text"
    except ImportError:
        return "Image (OCR not available)"


def load_text_file(path: Path) -> tuple[str, List[ExtractedImage]]:
    with path.open("r", encoding="utf-8", errors="ignore") as f:
        return f.read(), []


def load_pdf(path: Path) -> tuple[str, List[ExtractedImage]]:
    text_parts = []
    extracted_images = []
    reader = PdfReader(str(path))
    
    # First, try standard PDF object extraction
    for page_num, page in enumerate(reader.pages):
        page_text = page.extract_text() or ""
        if page_text:
            text_parts.append(page_text)
        
        # Extract images from PDF
        try:
            if "/Resources" in page and "/XObject" in page["/Resources"]:
                xobj = page["/Resources"]["/XObject"].get_object()
                for obj_name in xobj:
                    try:
                        obj = xobj[obj_name].get_object()
                        if obj.get("/Subtype") == "/Image":
                            try:
                                data = obj.get_data()
                                image = Image.open(io.BytesIO(data))
                                if image.size[0] > 50 and image.size[1] > 50:  # Filter out tiny images
                                    extracted_images.append(ExtractedImage(image, "pdf", page_num + 1))
                            except Exception as e:
                                logging.debug("Failed to decode image on page %d: %s", page_num + 1, e)
                    except Exception as e:
                        logging.debug("Error processing XObject on page %d: %s", page_num + 1, e)
        except Exception as e:
            logging.debug("Error accessing resources on page %d: %s", page_num + 1, e)
    
    # If we found no images with standard extraction, try pdf2image for screenshot/scanned PDFs
    if not extracted_images:
        try:
            from pdf2image import convert_from_path
            logging.info("Using pdf2image to render PDF pages as images (for screenshot-based PDFs)")
            images = convert_from_path(str(path), first_page=1, last_page=len(reader.pages))
            for page_num, img in enumerate(images, 1):
                if img.size[0] > 50 and img.size[1] > 50:
                    extracted_images.append(ExtractedImage(img, "pdf", page_num))
        except ImportError:
            logging.debug("pdf2image not available; cannot render PDF pages as images")
        except Exception as e:
            logging.warning("Failed to render PDF with pdf2image: %s (requires poppler-utils)", e)
    
    return "\n\n".join(text_parts), extracted_images


def load_docx(path: Path) -> tuple[str, List[ExtractedImage]]:
    doc = DocxDocument(str(path))
    paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
    extracted_images = []
    
    for inline_shape in doc.inline_shapes:
        try:
            r_id = inline_shape._inline.graphic.graphicData.pic.blipFill.blip.embed
            part = doc.part.related_parts[r_id]
            img_data = part.blob
            image = Image.open(io.BytesIO(img_data))
            if image.size[0] > 50 and image.size[1] > 50:  # Filter out tiny images
                extracted_images.append(ExtractedImage(image, "docx"))
        except Exception:
            pass
    
    for table in doc.tables:
        for row in table.rows:
            row_text = " \t ".join(cell.text for cell in row.cells if cell.text.strip())
            if row_text:
                paragraphs.append(row_text)
    
    return "\n\n".join(paragraphs), extracted_images


def load_pptx(path: Path) -> tuple[str, List[ExtractedImage]]:
    prs = PptxPresentation(str(path))
    text_parts = []
    extracted_images = []
    
    for slide_num, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text_parts.append(shape.text)
            if hasattr(shape, "image"):
                try:
                    img_data = shape.image.blob
                    image = Image.open(io.BytesIO(img_data))
                    if image.size[0] > 50 and image.size[1] > 50:  # Filter out tiny images
                        extracted_images.append(ExtractedImage(image, "pptx", slide_num + 1))
                except Exception:
                    pass
        
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame is not None:
            notes = slide.notes_slide.notes_text_frame.text
            if notes:
                text_parts.append(notes)
    
    return "\n\n".join([part for part in text_parts if part.strip()]), extracted_images


def load_image(path: Path) -> tuple[str, List[ExtractedImage]]:
    try:
        image = Image.open(str(path))
        if image.size[0] > 50 and image.size[1] > 50:  # Only include if reasonably sized
            return f"[Image file: {path.name}]", [ExtractedImage(image, "image", path.name)]
        else:
            return f"[Image file: {path.name} - too small]", []
    except Exception:
        return f"[Image file: {path.name} - could not load]", []


def load_document(path: str) -> tuple[str, List[ExtractedImage]]:
    path_obj = Path(path)
    if path_obj.is_dir():
        documents = []
        extracted_images = []
        for child in sorted(path_obj.iterdir()):
            if child.is_file():
                doc_text, imgs = load_document(str(child))
                documents.append(doc_text)
                extracted_images.extend(imgs)
        return "\n\n".join(documents), extracted_images

    ext = path_obj.suffix.lower()
    if ext in {".txt", ".md", ".csv", ".json"}:
        return load_text_file(path_obj)
    if ext == ".pdf":
        return load_pdf(path_obj)
    if ext == ".docx":
        return load_docx(path_obj)
    if ext == ".pptx":
        return load_pptx(path_obj)
    if ext in {".png", ".jpg", ".jpeg", ".bmp", ".gif", ".webp"}:
        return load_image(path_obj)

    raise ValueError(f"Unsupported input file type: {ext}")
